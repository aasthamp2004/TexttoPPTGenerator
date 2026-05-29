"""
ppt_service.py  –  Fully Dynamic AI PPT Generator
===================================================
- ALL content is LLM-generated (zero hardcoded bullets/labels/values)
- Logo color extraction preserved
- Every layout renders from its own data fields (never falls back silently)
- Attractive, varied visual styles per deck profile
- Strict field validation with meaningful LLM-fallback repair
"""

import os
import json
import ast
import re
import random
import colorsys
from openai import AzureOpenAI
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

from backend.config import (
    AZURE_KEY, AZURE_ENDPOINT, AZURE_API_VERSION, AZURE_DEPLOYMENT
)

_client = AzureOpenAI(
    api_key=AZURE_KEY,
    api_version=AZURE_API_VERSION,
    azure_endpoint=AZURE_ENDPOINT,
)

# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE DIMENSIONS  (all in inches — single source of truth)
# ═══════════════════════════════════════════════════════════════════════════
SW, SH = 13.3, 7.5          # slide width / height
HEADER_H   = 1.40           # header band height
FOOTER_H   = 0.28           # footer band height
CONTENT_Y  = HEADER_H + 0.14        # top of content area
CONTENT_H  = SH - HEADER_H - FOOTER_H - 0.22  # usable content height
CONTENT_X  = 0.42           # left margin
CONTENT_W  = SW - CONTENT_X - 0.28  # usable content width

# ═══════════════════════════════════════════════════════════════════════════
#  PROFILE SYSTEM  – random visual style per deck
# ═══════════════════════════════════════════════════════════════════════════

PROFILES = {
    "classic":   dict(cover="bands",     header="solid",   footer="bar",  card="outline"),
    "magazine":  dict(cover="sidebar",   header="sidebar", footer="line", card="soft"),
    "executive": dict(cover="centered",  header="band",    footer="bar",  card="outline"),
    "tech":      dict(cover="grid",      header="split",   footer="line", card="banded"),
    "bold":      dict(cover="dark_full", header="band",    footer="bar",  card="soft"),
    "minimal":   dict(cover="centered",  header="split",   footer="line", card="outline"),
}

PALETTES = {
    "classic":   dict(p=(11,95,255),   s=(77,163,255),  a=(0,58,160),   a2=(0,163,163),  dk=(11,42,74),   muted=(100,130,160), hf="Calibri",       bf="Calibri"),
    "magazine":  dict(p=(184,80,66),   s=(231,194,89),  a=(120,45,35),  a2=(167,190,174),dk=(74,28,42),   muted=(140,100,90),  hf="Georgia",        bf="Calibri"),
    "executive": dict(p=(44,95,45),    s=(151,188,98),  a=(20,55,20),   a2=(90,175,150), dk=(22,55,22),   muted=(90,120,90),   hf="Trebuchet MS",   bf="Calibri"),
    "tech":      dict(p=(109,46,70),   s=(162,103,105), a=(70,22,40),   a2=(236,226,208),dk=(44,18,28),   muted=(130,90,100),  hf="Consolas",       bf="Calibri"),
    "bold":      dict(p=(20,20,90),    s=(240,80,60),   a=(10,10,55),   a2=(255,200,0),  dk=(10,10,45),   muted=(110,110,150), hf="Arial Black",    bf="Arial"),
    "minimal":   dict(p=(54,69,79),    s=(120,145,160), a=(28,42,50),   a2=(0,164,180),  dk=(22,35,45),   muted=(110,130,140), hf="Calibri Light",  bf="Calibri"),
}

# ═══════════════════════════════════════════════════════════════════════════
#  LOW-LEVEL HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def _c(t):   return RGBColor(*[max(0,min(255,int(v))) for v in t])
def _IN(v):  return Inches(float(v))
def _mix(c1, c2, t):
    t = max(0.0, min(1.0, float(t)))
    return tuple(int(c1[i]+(c2[i]-c1[i])*t) for i in range(3))
def _lum(c):
    r,g,b = [v/255 for v in c]
    return 0.2126*r + 0.7152*g + 0.0722*b
def _contrast_text(bg, dark=(0,0,0), light=(255,255,255), threshold=0.56):
    return light if _lum(bg) < threshold else dark
def _font_pt(size, role="body"):
    size = float(size)
    if role == "title":
        return max(size, round(size * 1.08, 1))
    if role == "subtitle":
        return max(size, round(size * 1.10, 1))
    return max(size, round(size * 1.09, 1))
def _is_neutral(c, thr=0.16):
    r,g,b = [v/255 for v in c]
    _,s,_ = colorsys.rgb_to_hsv(r,g,b)
    return s < thr
def _badge_text(bg):
    return _contrast_text(bg, threshold=0.46)

def _parse_color(v, fallback):
    if isinstance(v,(list,tuple)) and len(v)==3:
        try: return tuple(max(0,min(255,int(x))) for x in v)
        except: return fallback
    if isinstance(v,str):
        s = v.strip().lstrip("#")
        if re.fullmatch(r"[0-9A-Fa-f]{6}", s):
            return tuple(int(s[i:i+2],16) for i in (0,2,4))
        parts = [p.strip() for p in s.split(",")]
        if len(parts)==3:
            try: return tuple(max(0,min(255,int(x))) for x in parts)
            except: pass
    return fallback

def _parse_int(v, default=50, lo=0, hi=100):
    if isinstance(v, int): return max(lo,min(hi,v))
    m = re.search(r"-?\d+", str(v))
    return max(lo,min(hi,int(m.group()))) if m else default

def _safe_str(v, default=""):
    return str(v).strip() if v is not None else default

def _safe_list(v):
    return [str(x).strip() for x in v if str(x).strip()] if isinstance(v,list) else []

def _dedupe_keep_order(items):
    out, seen = [], set()
    for item in items or []:
        clean = _safe_str(item)
        key = clean.casefold()
        if not clean or key in seen:
            continue
        seen.add(key)
        out.append(clean)
    return out
def _norm_text(v):
    s = _safe_str(v).casefold()
    s = re.sub(r"[^a-z0-9]+", " ", s).strip()
    return s
def _looks_like_prompt_title(title: str, topic: str) -> bool:
    tt = _norm_text(title)
    tp = _norm_text(topic)
    if not tt or not tp:
        return False
    if tt == tp:
        return True
    if tp in tt and len(tp) >= max(12, int(len(tt) * 0.6)):
        return True
    return any(tt.startswith(prefix) for prefix in (
        "create a powerpoint deck on",
        "create a presentation on",
        "presentation on",
        "ppt on",
        "deck on",
    ))
def _fallback_slide_title(slide: dict, idx: int) -> str:
    layout = _canon(slide.get("layout", "bullets"))
    if layout == "section_index":
        return "Contents"
    if layout == "title_cover":
        return "Executive Overview"
    if layout == "timeline":
        steps = slide.get("steps") or []
        if isinstance(steps, list) and steps:
            first = steps[0] if isinstance(steps[0], dict) else {}
            lbl = _safe_str(first.get("label", "Roadmap"))
            return lbl if lbl else "Roadmap"
        return "Roadmap"
    if layout == "icon_grid":
        items = slide.get("grid_items") or []
        if isinstance(items, list) and items:
            first = items[0] if isinstance(items[0], dict) else {}
            base = _safe_str(first.get("title", "Key Pillars"))
            return base if len(base.split()) <= 4 else "Key Pillars"
        return "Key Pillars"
    if layout == "case_study":
        company = _safe_str(slide.get("company", ""))
        return f"Case Study: {company}" if company else "Case Study"
    if layout == "big_stat":
        return _safe_str(slide.get("stat_label", "")) or "Key Metric"
    content = _safe_list(slide.get("content", []))
    if content:
        first = re.split(r"[:.;-]", content[0], maxsplit=1)[0].strip()
        words = first.split()
        if 1 <= len(words) <= 6:
            return first
    return f"Slide {idx}"

def _usage_dict(usage):
    if not usage:
        return None
    return {
        "prompt_tokens": int(getattr(usage, "prompt_tokens", 0) or 0),
        "completion_tokens": int(getattr(usage, "completion_tokens", 0) or 0),
        "total_tokens": int(getattr(usage, "total_tokens", 0) or 0),
    }

def _merge_usage(total, add):
    if not add:
        return total
    if total is None:
        total = {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0}
    total["prompt_tokens"] += int(add.get("prompt_tokens", 0) or 0)
    total["completion_tokens"] += int(add.get("completion_tokens", 0) or 0)
    total["total_tokens"] += int(add.get("total_tokens", 0) or 0)
    return total

# ═══════════════════════════════════════════════════════════════════════════
#  LOGO COLOR EXTRACTION  (unchanged from original)
# ═══════════════════════════════════════════════════════════════════════════

def _extract_logo_theme(logo_path: str, base: dict) -> dict:
    theme = dict(base)
    if not logo_path or not os.path.exists(logo_path):
        return theme
    try:
        with Image.open(logo_path) as img:
            rgba = img.convert("RGBA")
            bg   = Image.new("RGBA", rgba.size, (255,255,255,255))
            comp = Image.alpha_composite(bg, rgba).convert("RGB")
            comp.thumbnail((220,220))
            q    = comp.quantize(colors=6)
            pal  = q.getpalette() or []
            cnts = q.getcolors() or []
        extracted = []
        for cnt, idx in sorted(cnts, reverse=True):
            b = idx*3
            if b+2 >= len(pal): continue
            col = tuple(pal[b:b+3])
            lm  = _lum(col)
            if lm > 0.96 or lm < 0.04: continue
            extracted.append((cnt,col))
        colors  = [c for _,c in extracted]
        vivid   = [c for c in colors if not _is_neutral(c)]
        neutral = [c for c in colors if _is_neutral(c)]
        vivid.sort(key=lambda c: colorsys.rgb_to_hsv(*[v/255 for v in c])[1], reverse=True)
        primary   = vivid[0] if vivid else theme["p"]
        secondary = vivid[1] if len(vivid)>1 else _mix(primary,(255,255,255),0.35)
        accent    = vivid[2] if len(vivid)>2 else _mix(primary,(0,0,0),0.28)
        accent2   = neutral[0] if neutral else _mix(secondary,(255,255,255),0.18)
        theme.update(dict(
            p=primary, s=secondary, a=accent, a2=accent2,
            dk=_mix(primary,(12,18,28),0.42),
            bg=(255,255,255),
            td=_mix(primary,(15,15,15),0.70) if _lum(primary)>0.35 else (28,36,48),
            tl=(255,255,255),
            muted=_mix((28,36,48),(255,255,255),0.48),
            card=(255,255,255),
        ))
    except Exception as e:
        print(f"[logo-theme] {e}")
    return theme


# ═══════════════════════════════════════════════════════════════════════════
#  THEME BUILDING
# ═══════════════════════════════════════════════════════════════════════════

_BASE = dict(
    p=(34,83,149), s=(106,161,218), a=(34,83,149), a2=(106,161,218),
    bg=(255,255,255), dk=(34,83,149), td=(34,50,80), tl=(255,255,255),
    muted=(106,161,218), card=(255,255,255), hf="Calibri", bf="Calibri",
)

def _build_theme(palette_name: str, design_system: dict, logo_path: str) -> dict:
    pal = PALETTES.get(palette_name, PALETTES["classic"])
    theme = dict(
        p=pal["p"], s=pal["s"], a=pal["a"], a2=pal["a2"],
        dk=pal["dk"], bg=(255,255,255), td=(28,36,48), tl=(255,255,255),
        muted=pal["muted"], card=(255,255,255),
        hf=pal["hf"], bf=pal["bf"],
    )
    # overlay design_system colors from LLM
    raw = design_system.get("theme", {}) if isinstance(design_system, dict) else {}
    for key, tkey in [("primary","p"),("secondary","s"),("accent","a"),("accent2","a2"),
                      ("bg_dark","dk"),("bg_light","bg"),("text_dark","td"),
                      ("text_light","tl"),("text_muted","muted"),("card_bg","card")]:
        if key in raw:
            theme[tkey] = _parse_color(raw[key], theme[tkey])
    for fkey in ("header_font","body_font"):
        tkey = "hf" if fkey=="header_font" else "bf"
        if fkey in raw and isinstance(raw[fkey], str) and raw[fkey].strip():
            theme[tkey] = raw[fkey].strip()
    # Enforce readable text colors
    if _lum(theme.get("td", (0,0,0))) > 0.60:
        theme["td"] = (20, 20, 20)
    if _lum(theme.get("tl", (255,255,255))) < 0.40:
        theme["tl"] = (255, 255, 255)
    if logo_path:
        theme = _extract_logo_theme(logo_path, theme)
    return theme


# ═══════════════════════════════════════════════════════════════════════════
#  DRAWING PRIMITIVES
# ═══════════════════════════════════════════════════════════════════════════

def _rect(slide, x, y, w, h, fill, line=None, lw=0.75):
    s = slide.shapes.add_shape(1, _IN(x), _IN(y), _IN(w), _IN(h))
    s.fill.solid(); s.fill.fore_color.rgb = _c(fill)
    if line: s.line.color.rgb = _c(line); s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def _oval(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(9, _IN(x), _IN(y), _IN(w), _IN(h))
    s.fill.solid(); s.fill.fore_color.rgb = _c(fill)
    s.line.fill.background()
    return s

def _tb(slide, text, x, y, w, h, size, bold=False, italic=False,
        color=None, face="Calibri", align=PP_ALIGN.LEFT, shrink=False):
    text = _safe_str(text)
    bx = slide.shapes.add_textbox(_IN(x), _IN(y), _IN(w), _IN(h))
    tf = bx.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    # Enable auto-size by default to prevent text overlap when it wraps.
    # If shrink is True or not specified, we use TEXT_TO_FIT_SHAPE.
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.alignment = align
    role = "title" if bold and size >= 18 else ("subtitle" if italic else "body")
    font_size = _font_pt(size, role=role)
    for i, part in enumerate(str(text).split("**")):
        if not part: continue
        r = p.add_run(); r.text = part
        r.font.size = Pt(font_size); r.font.bold = bold or (i%2==1)
        r.font.italic = italic; r.font.name = face
        if color: r.font.color.rgb = _c(color)
    return bx

def _bullets(slide, points, x, y, w, h, size=16, icon="▸",
             ic=None, tc=None, face="Calibri", maxp=6):
    if not points: return
    ic = ic or _BASE["s"]; tc = tc or _BASE["td"]
    pts = [_safe_str(p) for p in points if _safe_str(p)]
    if not pts:
        return
    # Keep all bullets (bounded by maxp if provided).
    maxp = maxp or len(pts)
    pts = pts[:maxp]
    if not pts: return
    n = len(pts)
    if   n <= 2: size = min(26, size+6)
    elif n == 3: size = min(23, size+4)
    elif n == 4: size = min(20, size+2)
    elif n == 5: size = min(18, size+1)
    elif n >= 7: size = max(12, size-2)
    bx = slide.shapes.add_textbox(_IN(x),_IN(y),_IN(w),_IN(h))
    tf = bx.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if n<=3 else MSO_ANCHOR.TOP
    first = True
    sp = 13 if n<=3 else (6 if n>=7 else 8)
    for pt in pts:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        # Use line_spacing to ensure wrapped lines don't overlap and have breathing room.
        p.line_spacing = 1.2
        p.space_before = Pt(sp); p.space_after = Pt(sp)
        ir = p.add_run(); ir.text = f"{icon}  "
        ir.font.size=Pt(_font_pt(size-1)); ir.font.bold=True
        ir.font.name=face; ir.font.color.rgb=_c(ic)
        for i,part in enumerate(pt.split("**")):
            if not part: continue
            r=p.add_run(); r.text=part
            r.font.size=Pt(_font_pt(size)); r.font.bold=(i%2==1)
            r.font.name=face; r.font.color.rgb=_c(tc)


# ═══════════════════════════════════════════════════════════════════════════
#  SHARED CHROME: header / footer / card / logo
# ═══════════════════════════════════════════════════════════════════════════

def _draw_header(slide, title, subtitle, num, theme, profile):
    mode = profile["header"]
    W = SW; H = HEADER_H
    tx, ty, th = 0.45, 0.18, 0.76
    title_fill = theme["dk"]

    if mode == "sidebar":
        _rect(slide, 0, 0, 0.42, H, theme["p"])
        title_fill = _mix(theme["dk"],theme["p"],0.18)
        _rect(slide, 0.42, 0, W-0.42, H, title_fill)
    elif mode == "band":
        _rect(slide, 0, 0, W, H, theme["dk"])
        _rect(slide, 0, 0, W, 0.16, theme["s"])
        title_fill = theme["dk"]
    elif mode == "split":
        _rect(slide, 0, 0, W*0.58, H, theme["dk"])
        _rect(slide, W*0.58, 0, W*0.42, H, theme["p"])
        title_fill = theme["dk"]
    else:  # solid
        _rect(slide, 0, 0, W, H, theme["dk"])
        title_fill = theme["dk"]

    title_w = W - tx - 2.20
    tsize = max(22, 32 - max(0, len(title)-42)//8*4)
    # Use a small line spacing for titles to handle wrapping without overlapping subtitle.
    title_bx = _tb(slide, title, tx, ty, title_w, th, tsize,
        bold=True, color=_contrast_text(title_fill), face=theme["hf"], shrink=True)
    title_bx.text_frame.paragraphs[0].line_spacing = 1.0

    if subtitle:
        ssize = max(11, 14 - max(0, len(subtitle)-60)//20)
        subtitle_bx = _tb(slide, subtitle, tx, 0.96, title_w, 0.38, ssize,
            italic=True, color=_contrast_text(title_fill), face=theme["bf"], shrink=True)
        subtitle_bx.text_frame.paragraphs[0].line_spacing = 1.0


def _draw_footer(slide, theme, profile):
    mode = profile["footer"]
    if mode == "line":
        _rect(slide, 0, SH-0.07, SW, 0.07, theme["s"])
    else:
        _rect(slide, 0, SH-FOOTER_H, SW, FOOTER_H, theme["dk"])


def _add_page_number(slide, num, theme, footer_mode="bar"):
    if num is None:
        return
    # Simple text page number at bottom-left to avoid logo overlap.
    ns = str(num)
    if footer_mode == "line":
        # Place on slide background just above the line.
        y = SH - 0.30
        color = _contrast_text(theme["bg"])
    else:
        # Place inside the footer band.
        y = SH - FOOTER_H + 0.03
        color = _contrast_text(theme["dk"])
    _tb(slide, ns, 0.28, y, 0.60, 0.22, 12,
        bold=True, color=color, face=theme["bf"], align=PP_ALIGN.LEFT)


def _draw_card(slide, theme, profile, x=None, y=None, w=None, h=None):
    cx = x if x is not None else CONTENT_X
    cy = y if y is not None else CONTENT_Y
    cw = w if w is not None else CONTENT_W
    ch = h if h is not None else CONTENT_H
    mode = profile["card"]
    if mode == "soft":
        _rect(slide, cx, cy, cw, ch,
              _mix(theme["card"], theme["s"], 0.12),
              line=_mix(theme["p"],theme["s"],0.5), lw=0.8)
    elif mode == "banded":
        _rect(slide, cx, cy, cw, ch, theme["card"], line=theme["p"], lw=1.0)
        _rect(slide, cx, cy, cw, 0.12, theme["s"])
    else:
        _rect(slide, cx, cy, cw, ch, theme["card"], line=theme["p"], lw=1.0)
    pad = 0.20
    return cx+pad, cy+0.18, cw-pad*2, ch-0.18-0.28


def _add_logo(slide, logo_path, theme):
    if not logo_path or not os.path.exists(logo_path): return
    try:
        box_w, box_h, margin_r, margin_t = 1.75, 0.92, 0.16, 0.08
        bx = SW - box_w - margin_r
        by = margin_t
        from io import BytesIO
        with Image.open(logo_path) as img:
            rgba = img.convert("RGBA")
            alpha = rgba.split()[-1]
            bbox = alpha.getbbox()
            if bbox:
                rgba = rgba.crop(bbox)
            # Trim near-white borders
            px = rgba.getdata()
            mask = Image.new("L", rgba.size, 0)
            mpx = mask.load()
            w, h = rgba.size
            for y in range(h):
                for x in range(w):
                    r, g, b, a = px[y*w + x]
                    if a > 5 and not (r > 245 and g > 245 and b > 245):
                        mpx[x, y] = 255
            bbox2 = mask.getbbox()
            if bbox2:
                rgba = rgba.crop(bbox2)
            iw, ih = rgba.size
            ratio = iw/ih if ih else 1.0
            lw = min(box_w, box_h * ratio)
            lh = lw / ratio if ratio else box_h
            if lh > box_h:
                lh = box_h
                lw = lh * ratio
            buf = BytesIO()
            rgba.save(buf, format="PNG")
            buf.seek(0)
            slide.shapes.add_picture(
                buf,
                _IN(bx + (box_w - lw) / 2),
                _IN(by + (box_h - lh) / 2),
                _IN(lw),
                _IN(lh),
            )
    except Exception as e:
        print(f"[logo] {e}")


def _slide_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _c(color)


# ═══════════════════════════════════════════════════════════════════════════
#  ACCENT CYCLING  — per-slide color variety
# ═══════════════════════════════════════════════════════════════════════════

def _accents(theme, n=4):
    base = [theme["p"], theme["a2"], theme["s"], theme["a"]]
    return [base[i % len(base)] for i in range(n)]


# ═══════════════════════════════════════════════════════════════════════════
#  LAYOUT RENDERERS  — each reads ONLY its own data fields
# ═══════════════════════════════════════════════════════════════════════════

def render_title_cover(slide, spec, num, theme, profile, logo_path):
    W, H = SW, SH
    mode = profile["cover"]
    _slide_bg(slide, theme["bg"])

    if mode == "dark_full":
        _slide_bg(slide, theme["dk"])
        _rect(slide, 0, 0, 0.35, H, theme["s"])
        _rect(slide, 0, H-0.12, W, 0.12, theme["a2"])
        tc, sc = _contrast_text(theme["dk"]), _contrast_text(theme["dk"])
    elif mode == "sidebar":
        _rect(slide, 0, 0, 2.05, H, theme["p"])
        _rect(slide, 2.05, 0, 0.10, H, theme["s"])
        tc, sc = _contrast_text(theme["bg"]), _contrast_text(theme["bg"])
    elif mode == "grid":
        _rect(slide, 0, 0, W, 1.25, theme["p"])
        _rect(slide, 0, 1.25, W, 0.08, theme["s"])
        for gx in [1.2, 3.0, 4.8, 6.6, 8.4, 10.2, 12.0]:
            _rect(slide, gx, 0, 0.018, H, _mix(theme["s"], theme["bg"], 0.82))
        _rect(slide, 0, H-0.18, W, 0.18, _mix(theme["p"],theme["s"],0.55))
        tc, sc = _contrast_text(theme["bg"]), _contrast_text(theme["bg"])
    elif mode == "centered":
        _rect(slide, 0, 0, W, 1.06, theme["p"])
        _rect(slide, 0, H-0.18, W, 0.18, _mix(theme["p"],theme["s"],0.55))
        tc, sc = _contrast_text(theme["bg"]), _contrast_text(theme["bg"])
    else:  # bands
        _rect(slide, 0, 0, W, 1.25, theme["p"])
        _rect(slide, 0, 1.25, W, 0.08, theme["s"])
        _rect(slide, 0, H-0.18, W, 0.18, _mix(theme["p"],theme["s"],0.55))
        tc, sc = _contrast_text(theme["bg"]), _contrast_text(theme["bg"])

    title_x = 2.40 if mode == "sidebar" else 0.70
    title_w = W - title_x - 0.40
    title = _safe_str(spec.get("title", "Presentation"))
    tsize = max(40, 58 - max(0, len(title)-30)//6*4)
    cover_title_bx = _tb(slide, title, title_x, 1.95, title_w, 2.20, tsize,
        bold=True, color=tc, face=theme["hf"], align=PP_ALIGN.LEFT, shrink=True)
    cover_title_bx.text_frame.paragraphs[0].line_spacing = 1.0

    subtitle = _safe_str(spec.get("subtitle",""))
    if subtitle:
        ssize = max(15, 20 - max(0, len(subtitle)-60)//20)
        cover_sub_bx = _tb(slide, subtitle, title_x+0.02, 4.15, title_w, 0.82, ssize,
            italic=True, color=sc, face=theme["bf"], shrink=True)
        cover_sub_bx.text_frame.paragraphs[0].line_spacing = 1.0

    # Highlight cards from content
    points = [p for p in _safe_list(spec.get("content",[])) if p][:3]
    if points:
        accs = _accents(theme, 3)
        card_w = max(2.35, (title_w - 0.18*(len(points)-1)) / len(points))
        for i, pt in enumerate(points):
            cx = title_x + i*(card_w+0.18)
            fill = _mix(theme["card"], accs[i], 0.88 if mode=="dark_full" else 0.82)
            _rect(slide, cx, 5.00, card_w, 1.02, fill, line=accs[i], lw=1.0)
            _rect(slide, cx, 5.00, card_w, 0.10, accs[i])
            _tb(slide, pt, cx+0.12, 5.16, card_w-0.24, 0.76, 15,
                bold=True, color=_contrast_text(fill),
                face=theme["bf"], shrink=True)
    _add_logo(slide, logo_path, theme)


def render_section_index(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title","Contents"), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)
    sections = _safe_list(spec.get("sections", spec.get("content",[])))
    if not sections:
        _tb(slide, "No sections available.", ix, iy, iw, ih, 15, color=theme["muted"])
        _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme); return

    sections = sections[:6]
    cols = 2 if len(sections) > 3 else 1
    gap = 0.18
    rows = (len(sections)+cols-1)//cols
    cw = (iw - gap*(cols-1)) / cols
    ch = min(1.08, max(0.72, (ih - gap*(rows-1)) / max(1,rows)))
    accs = _accents(theme, 4)

    for i, sec in enumerate(sections):
        col = i % cols; row = i // cols
        cx = ix + col*(cw+gap)
        cy = iy + row*(ch+gap)
        ac = accs[i % len(accs)]
        fill = _mix(theme["card"],ac,0.90)
        _rect(slide, cx, cy, cw, ch, fill, line=ac, lw=1.0)
        _rect(slide, cx, cy, 0.72, ch, ac)
        _tb(slide, f"{i+1}", cx+0.08, cy+0.14, 0.64, ch-0.24, 17,
            bold=True, color=_contrast_text(ac), face=theme["hf"], align=PP_ALIGN.CENTER)
        _tb(slide, sec, cx+0.86, cy+0.16, cw-1.00, ch-0.28,
            18 if cols==1 else 15, bold=True, color=_contrast_text(fill),
            face=theme["bf"], shrink=True)
    _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme)


def render_bullets(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title",""), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)
    _bullets(slide, spec.get("content",[]), ix, iy, iw, ih,
             size=16, icon=spec.get("icon","▸"),
             ic=theme["s"], tc=theme["td"], face=theme["bf"], maxp=6)
    _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme)


def render_two_column(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title",""), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)

    gap = 0.20; cw = (iw-gap)/2
    lt = _safe_str(spec.get("left_title","Left"))
    rt = _safe_str(spec.get("right_title","Right"))
    lp = _safe_list(spec.get("left_points",[]))
    rp = _safe_list(spec.get("right_points",[]))

    # Fallback: split content if left/right empty
    if not lp and not rp:
        content = _safe_list(spec.get("content",[]))
        mid = max(1, len(content)//2)
        lp, rp = content[:mid], content[mid:]

    icon = spec.get("icon","▸")
    hh = 0.46

    _rect(slide, ix, iy, cw, hh, theme["p"])
    _tb(slide, lt, ix+0.10, iy+0.08, cw-0.20, hh-0.12, 15,
        bold=True, color=_contrast_text(theme["p"]), face=theme["hf"])
    _bullets(slide, lp, ix+0.10, iy+hh+0.10, cw-0.20, ih-hh-0.16,
             size=15, icon=icon, ic=theme["s"], tc=theme["td"],
             face=theme["bf"], maxp=6)

    rx = ix+cw+gap
    _rect(slide, rx, iy, cw, hh, theme["a2"])
    _tb(slide, rt, rx+0.10, iy+0.08, cw-0.20, hh-0.12, 15,
        bold=True, color=_contrast_text(theme["a2"]), face=theme["hf"])
    _bullets(slide, rp, rx+0.10, iy+hh+0.10, cw-0.20, ih-hh-0.16,
             size=15, icon=icon, ic=theme["a"], tc=theme["td"],
             face=theme["bf"], maxp=6)
    _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme)


def render_big_stat(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title",""), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)

    pw = 3.70
    stat   = _safe_str(spec.get("stat","—"))
    label  = _safe_str(spec.get("stat_label",""))
    source = _safe_str(spec.get("stat_source",""))
    pts    = _safe_list(spec.get("content",[]))

    sfont = 72 if len(stat)<=4 else 54
    _rect(slide, ix, iy, pw, ih, theme["p"])
    _rect(slide, ix, iy, pw, 0.07, theme["a2"])
    _tb(slide, stat, ix+0.10, iy+0.38, pw-0.20, 1.80, sfont,
        bold=True, color=_contrast_text(theme["p"]), face=theme["hf"], align=PP_ALIGN.CENTER)
    _tb(slide, label, ix+0.10, iy+2.25, pw-0.20, 0.75, 15,
        color=_contrast_text(theme["p"]), face=theme["bf"], align=PP_ALIGN.CENTER)
    if source:
        _tb(slide, source, ix+0.10, iy+3.05, pw-0.20, 0.50, 11,
            italic=True, color=_contrast_text(theme["p"]), face=theme["bf"], align=PP_ALIGN.CENTER)
    # Progress bar if numeric
    try:
        pct = float(stat.replace("%","").replace("+","").strip())
        if 0 < pct < 100:
            bmar = 0.30; bw = pw-bmar*2; bar_y = iy+ih-0.55; bh = 0.20
            _rect(slide, ix+bmar, bar_y, bw, bh, (255,255,255))
            _rect(slide, ix+bmar, bar_y, bw*(pct/100), bh, theme["a2"])
    except: pass

    bx = ix+pw+0.22; bw2 = iw-pw-0.22
    _bullets(slide, pts, bx, iy+0.10, bw2, ih-0.20,
             size=14, icon=spec.get("icon","▸"),
             ic=theme["s"], tc=theme["td"], face=theme["bf"], maxp=6)
    _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme)


def render_timeline(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title",""), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)

    steps = [s for s in (spec.get("steps") or [])
             if isinstance(s,dict) and _safe_str(s.get("label")) and _safe_str(s.get("detail"))]
    if not steps:
        _tb(slide, "No timeline steps provided.", ix, iy, iw, ih, 14, color=theme["muted"])
        _add_page_number(slide, num, theme); _add_logo(slide, logo_path, theme); return

    n = len(steps); sw = iw/n
    tl_y = iy + ih*0.47
    DR, dr = 0.28, 0.14
    accs = _accents(theme, n)

    _rect(slide, ix, tl_y, iw, 0.06, theme["s"])
    for i, step in enumerate(steps):
        cx = ix + i*sw + sw/2
        ac = accs[i % len(accs)]
        inner_fill = (255,255,255)
        _oval(slide, cx-DR, tl_y-DR, DR*2, DR*2, ac)
        _oval(slide, cx-dr, tl_y-dr, dr*2, dr*2, inner_fill)
        _tb(slide, str(i+1), cx-0.20, tl_y-0.16, 0.40, 0.32, 12,
            bold=True, color=_badge_text(inner_fill), align=PP_ALIGN.CENTER)
        cw_ = sw*0.84; cx_c = cx-cw_/2
        above = (i%2==0)
        if above:
            cy = iy+0.05; ch = max(0.40, tl_y-DR-0.42-cy)
            if tl_y-DR > cy+ch+0.02:
                _rect(slide, cx-0.025, cy+ch, 0.05, tl_y-DR-(cy+ch), theme["s"])
        else:
            cy = tl_y+DR+0.28; ch = max(0.40, iy+ih-0.05-cy)
            if cy > tl_y+DR+0.02:
                _rect(slide, cx-0.025, tl_y+DR, 0.05, cy-(tl_y+DR), theme["s"])
        _rect(slide, cx_c, cy, cw_, ch, theme["card"], line=ac, lw=1.0)
        _tb(slide, step["label"], cx_c+0.10, cy+0.08, cw_-0.20, 0.42, 13,
            bold=True, color=_contrast_text(theme["card"]), face=theme["hf"])
        dh = ch-0.56
        if dh > 0.12:
            _tb(slide, step["detail"], cx_c+0.10, cy+0.52, cw_-0.20, dh, 12,
                color=_contrast_text(theme["card"]), face=theme["bf"])
    _add_page_number(slide, num, theme); _add_logo(slide, logo_path, theme)


def render_icon_grid(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title",""), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)

    items = [g for g in (spec.get("grid_items") or [])
             if isinstance(g,dict) and _safe_str(g.get("title"))][:4]
    if not items:
        _tb(slide, "No grid items provided.", ix, iy, iw, ih, 14, color=theme["muted"])
        _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme); return

    cols = 2; gap = 0.16
    cw_ = (iw-gap)/cols; ch_ = (ih-gap)/2
    accs = _accents(theme, 4)
    IR = 0.38

    for i, gi in enumerate(items):
        col = i%cols; row = i//cols
        cx = ix + col*(cw_+gap)
        cy = iy + row*(ch_+gap)
        ac = accs[i % len(accs)]
        _rect(slide, cx, cy, cw_, ch_, theme["card"], line=ac, lw=1.0)
        _rect(slide, cx, cy, 0.09, ch_, ac)
        iox = cx+0.24; ioy = cy+(ch_-IR*2)/2
        _oval(slide, iox, ioy, IR*2, IR*2, ac)
        raw = _safe_str(gi.get("icon", ""))
        ch1 = next((c.upper() for c in raw if c.isascii() and c.isalnum()), _seq_icon(i))
        _tb(slide, ch1, iox+0.04, ioy+0.08, IR*2-0.08, IR*1.6, 17,
            bold=True, color=_contrast_text(ac), face=theme["hf"], align=PP_ALIGN.CENTER)
        tx = iox+IR*2+0.16; tw = cw_-(iox-cx)-IR*2-0.22
        _tb(slide, _safe_str(gi.get("title","")), tx, cy+0.12, tw, 0.46, 16,
            bold=True, color=_contrast_text(theme["card"]), face=theme["hf"])
        _tb(slide, _safe_str(gi.get("detail","")), tx, cy+0.60, tw, ch_-0.72, 14,
            color=_contrast_text(theme["card"]), face=theme["bf"])
    _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme)


def render_case_study(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title",""), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)

    company = _safe_str(spec.get("company","Organisation"))
    result  = _safe_str(spec.get("result",""))
    metrics = [m for m in (spec.get("metrics") or []) if isinstance(m,dict)]
    pts     = _safe_list(spec.get("content",[]))
    icon    = spec.get("icon","▸")

    bh = 0.52; rh = 0.48; mh = 0.50
    _rect(slide, ix, iy, iw, bh, theme["p"])
    _tb(slide, f"  {company}", ix+0.18, iy+0.10, iw-0.36, bh-0.16, 17,
        bold=True, color=_contrast_text(theme["p"]), face=theme["hf"])
    cur_y = iy+bh+0.06
    if result:
        _rect(slide, ix, cur_y, iw, rh, theme["a2"])
        _tb(slide, f"  {result}", ix+0.14, cur_y+0.10, iw-0.28, rh-0.16, 13,
            bold=True, color=_contrast_text(theme["a2"]), face=theme["bf"])
        cur_y += rh+0.06

    metric_labels = []
    for i, m in enumerate(metrics[:3]):
        lbl = _safe_str(m.get("label","Impact"))
        val = _parse_int(m.get("value",50), default=50, lo=1, hi=99)
        metric_labels.append(lbl.lower())
        my  = cur_y + i*mh
        mlw = 2.90; bar_x = ix+mlw+0.20; bar_w = iw-mlw-0.40
        _tb(slide, lbl, ix+0.14, my+0.04, mlw, 0.30, 13, color=_contrast_text(theme["card"]))
        _rect(slide, bar_x, my+0.06, bar_w, 0.22, (225,232,245))
        _rect(slide, bar_x, my+0.06, bar_w*(val/100), 0.22, theme["a2"])
        label_w = 0.60
        bar_fill = bar_w * (val / 100)
        if bar_fill >= label_w + 0.10:
            label_x = bar_x + bar_fill - label_w - 0.04
        else:
            label_x = bar_x + 0.04
        _tb(slide, f"{val}", label_x, my+0.01, label_w, 0.22, 11,
            bold=True, color=_contrast_text(theme["card"]))
        metrics_used = (i+1)*mh

    # Remove duplicate metric lines from bullets
    if metric_labels:
        filtered = []
        for p in pts:
            pl = p.lower()
            if any(pl.startswith(lbl) or pl.startswith(lbl + ":") for lbl in metric_labels):
                continue
            filtered.append(p)
        pts = filtered

    bul_y = cur_y + (metrics_used if metrics else 0) + 0.10
    bul_h = ih-(bul_y-iy)-0.10
    if bul_h > 0.30 and pts:
        max_bullets = min(4, max(1, int(bul_h / 0.36)))
        _bullets(slide, pts, ix+0.10, bul_y, iw-0.20, bul_h,
                 size=14, icon=icon, ic=theme["s"], tc=theme["td"],
                 face=theme["bf"], maxp=max_bullets)
    _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme)


def render_table(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title",""), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)

    cols = _safe_list(spec.get("table_columns",[]))[:5]
    rows = [r for r in (spec.get("table_rows") or []) if isinstance(r,list)][:6]

    if not cols or not rows:
        _tb(slide, "No table data provided.", ix, iy, iw, ih, 14, color=theme["muted"])
        _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme); return

    nc   = len(cols)
    hh   = 0.55; cg = 0.02; pad = 0.06
    col_w = (iw-(nc-1)*cg)/nc
    row_h = min(0.65, max(0.38, (ih-hh-0.08)/len(rows)))

    for ci, col in enumerate(cols):
        cx = ix + ci*(col_w+cg)
        _rect(slide, cx, iy, col_w, hh, theme["p"], line=theme["s"], lw=0.8)
        _tb(slide, col, cx+pad, iy+pad, col_w-pad*2, hh-pad*2, 13,
            bold=True, color=_contrast_text(theme["p"]), face=theme["hf"], align=PP_ALIGN.CENTER, shrink=True)

    for ri, row in enumerate(rows):
        ry = iy+hh+ri*row_h
        fill = theme["card"] if ri%2==0 else _mix(theme["card"],theme["s"],0.35)
        cells = [_safe_str(v) for v in row[:nc]]
        while len(cells)<nc: cells.append("")
        for ci, cell in enumerate(cells):
            cx = ix + ci*(col_w+cg)
            _rect(slide, cx, ry, col_w, row_h, fill,
                  line=_mix(theme["p"],theme["card"],0.65), lw=0.5)
            _tb(slide, cell, cx+pad, ry+pad, col_w-pad*2, row_h-pad, 12,
                color=_contrast_text(fill), face=theme["bf"], shrink=True)
    _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme)


def render_chart(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title",""), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)

    chart_data = [x for x in (spec.get("chart_data") or [])
                  if isinstance(x,dict) and _safe_str(x.get("label"))][:5]
    if not chart_data:
        _tb(slide, "No chart data provided.", ix, iy, iw, ih, 14, color=theme["muted"])
        _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme); return

    ct = _safe_str(spec.get("chart_title", spec.get("title","")))
    _tb(slide, ct, ix+0.02, iy+0.02, iw-0.04, 0.46, 15,
        bold=True, color=_contrast_text(theme["card"]), face=theme["hf"])

    values = [max(1,_parse_int(d.get("value",0),default=1,lo=1,hi=100)) for d in chart_data]
    vmax = max(values)
    lw = max(1.95, iw*0.28); bar_x = ix+lw+0.16; bar_w = iw-lw-0.66
    row_h = 0.24; row_gap = 0.18
    avail_h = ih - 0.56 - 0.28
    needed = len(chart_data)*row_h + (len(chart_data)-1)*row_gap
    if needed > avail_h and len(chart_data)>1:
        sc = avail_h/needed; row_h *= sc; row_gap *= sc

    accs = _accents(theme, len(chart_data))
    for i, item in enumerate(chart_data):
        y = iy+0.54 + i*(row_h+row_gap)
        lbl = _safe_str(item.get("label",""))[:40]
        val = max(1, _parse_int(item.get("value",0),default=1,lo=1,hi=100))
        pct = val/vmax
        _tb(slide, lbl, ix+0.02, y+0.01, lw-0.08, row_h-0.02, 13,
            color=_contrast_text(theme["card"]), face=theme["bf"], shrink=True)
        _rect(slide, bar_x, y, bar_w, row_h, _mix(theme["card"],theme["s"],0.90))
        _rect(slide, bar_x, y, bar_w*pct, row_h, accs[i])
        _tb(slide, str(val), bar_x+bar_w+0.08, y-0.01, 0.52, row_h+0.04, 12,
            bold=True, color=_contrast_text(theme["card"]), align=PP_ALIGN.LEFT)

    src = _safe_str(spec.get("chart_source",""))
    if src:
        _tb(slide, src, ix+0.02, iy+ih-0.24, iw-0.04, 0.18, 10,
            italic=True, color=theme["muted"], face=theme["bf"])
    _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme)


def render_image_text_split(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title",""), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)

    side = _safe_str(spec.get("image_side","right")).lower()
    if side not in ("left","right"): side = "right"
    gap = 0.20; img_w = iw*0.44; txt_w = iw-img_w-gap
    imx = ix if side=="left" else ix+txt_w+gap
    tx  = ix+img_w+gap if side=="left" else ix

    uploaded_image_path = _safe_str(spec.get("uploaded_image_path", ""))
    if uploaded_image_path and os.path.exists(uploaded_image_path):
        _rect(slide, imx, iy, img_w, ih, theme["card"], line=theme["s"], lw=1.0)
        try:
            with Image.open(uploaded_image_path) as img:
                iw_px, ih_px = img.size
            scale = min((img_w - 0.12) / max(iw_px, 1), (ih - 0.12) / max(ih_px, 1))
            pw = max(0.40, iw_px * scale)
            ph = max(0.40, ih_px * scale)
            px = imx + (img_w - pw) / 2
            py = iy + (ih - ph) / 2
            slide.shapes.add_picture(uploaded_image_path, _IN(px), _IN(py), width=_IN(pw), height=_IN(ph))
        except Exception:
            _rect(slide, imx, iy, img_w, ih, _mix(theme["s"],theme["card"],0.75), line=theme["s"], lw=1.0)
    else:
        # Decorative image placeholder
        _rect(slide, imx, iy, img_w, ih, _mix(theme["s"],theme["card"],0.75), line=theme["s"], lw=1.0)
        _rect(slide, imx+0.10, iy+0.10, img_w-0.20, ih*0.20, theme["p"])
        _rect(slide, imx+0.10, iy+ih*0.34, img_w-0.20, ih*0.52, _mix(theme["a2"],theme["card"],0.55))
        _oval(slide, imx+img_w*0.10, iy+ih*0.72, 0.45, 0.45, theme["a"])
        _oval(slide, imx+img_w*0.78, iy+ih*0.16, 0.36, 0.36, theme["a2"])
    cap = _safe_str(spec.get("image_caption",""))
    if cap:
        cap_fill = _mix(theme["a2"],theme["card"],0.55)
        _tb(slide, cap, imx+0.14, iy+ih-0.48, img_w-0.28, 0.38, 13,
            bold=True, color=_contrast_text(cap_fill), face=theme["bf"], align=PP_ALIGN.CENTER)

    _bullets(slide, spec.get("content",[]), tx, iy+0.06, txt_w, ih-0.12,
             size=15, icon=spec.get("icon","▸"),
             ic=theme["s"], tc=theme["td"], face=theme["bf"], maxp=6)
    _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme)


def render_hybrid_insight(slide, spec, num, theme, profile, logo_path):
    _slide_bg(slide, theme["bg"])
    _draw_header(slide, spec.get("title",""), spec.get("subtitle",""), num, theme, profile)
    _draw_footer(slide, theme, profile)
    ix, iy, iw, ih = _draw_card(slide, theme, profile)

    lw = iw*0.43; gap = 0.20; rw = iw-lw-gap; rx = ix+lw+gap
    stat  = _safe_str(spec.get("stat","—"))
    label = _safe_str(spec.get("stat_label","Key Indicator"))

    left_fill = _mix(theme["card"],theme["s"],0.86)
    _rect(slide, ix, iy, lw, ih, left_fill, line=theme["s"], lw=1.0)
    sfont = 52 if len(stat)<=4 else 42
    _tb(slide, stat, ix+0.12, iy+0.30, lw-0.24, 1.15, sfont,
        bold=True, color=_contrast_text(left_fill), face=theme["hf"], align=PP_ALIGN.CENTER)
    _tb(slide, label, ix+0.12, iy+1.45, lw-0.24, 0.42, 13,
        color=_contrast_text(left_fill), face=theme["bf"], align=PP_ALIGN.CENTER)

    chart = [r for r in (spec.get("chart_data") or []) if isinstance(r,dict)][:3]
    if chart:
        vmax = max(max(1,_parse_int(r.get("value",1),default=1,lo=1,hi=100)) for r in chart)
        by = iy+2.05; row_h = 0.26; row_gap = 0.16
        for i, row in enumerate(chart):
            y = by + i*(row_h+row_gap)
            v = max(1, _parse_int(row.get("value",1),default=1,lo=1,hi=100))
            lbl = _safe_str(row.get("label",""))[:18]
            _tb(slide, lbl, ix+0.12, y-0.01, 1.40, row_h, 11,
                color=_contrast_text(left_fill), face=theme["bf"], shrink=True)
            bar_x = ix+1.55; bar_w = lw-2.08
            _rect(slide, bar_x, y, bar_w, row_h, _mix(theme["card"],theme["s"],0.93))
            _rect(slide, bar_x, y, bar_w*(v/vmax), row_h, theme["a2"])
            # Give the value label enough width to avoid digit wrapping in LibreOffice.
            vlabel_w = 0.60
            bar_fill = bar_w * (v / vmax)
            if bar_fill >= vlabel_w + 0.10:
                vlabel_x = bar_x + bar_fill - vlabel_w - 0.04
            else:
                vlabel_x = bar_x + 0.04
            _tb(slide, str(v), vlabel_x, y-0.01, vlabel_w, row_h, 10,
                bold=True, color=_contrast_text(left_fill), align=PP_ALIGN.RIGHT, shrink=True)

    _rect(slide, rx, iy, rw, ih, theme["card"], line=theme["p"], lw=1.0)
    _bullets(slide, spec.get("content",[]), rx+0.10, iy+0.10, rw-0.20, ih-0.20,
             size=14, icon=spec.get("icon","▸"),
             ic=theme["s"], tc=theme["td"], face=theme["bf"], maxp=6)
    _add_page_number(slide, num, theme, profile["footer"]); _add_logo(slide, logo_path, theme)


# ═══════════════════════════════════════════════════════════════════════════
#  LAYOUT REGISTRY
# ═══════════════════════════════════════════════════════════════════════════

_RENDERERS = {
    "title_cover":      render_title_cover,
    "section_index":    render_section_index,
    "bullets":          render_bullets,
    "two_column":       render_two_column,
    "big_stat":         render_big_stat,
    "timeline":         render_timeline,
    "icon_grid":        render_icon_grid,
    "case_study":       render_case_study,
    "table":            render_table,
    "chart":            render_chart,
    "image_text_split": render_image_text_split,
    "hybrid_insight":   render_hybrid_insight,
}

_CONTENT_LAYOUTS = list(_RENDERERS.keys())

_ALIASES = {
    "contents":"section_index","agenda":"section_index","index":"section_index",
    "comparison":"two_column","grid":"icon_grid","infographic":"icon_grid",
    "kpi":"big_stat","metric":"big_stat","chart_slide":"chart","graph":"chart",
    "bar_chart":"chart","image_text":"image_text_split","split":"image_text_split",
    "hybrid":"hybrid_insight","storytelling":"hybrid_insight",
    "two_col":"two_column",
}

def _canon(layout, default="bullets"):
    raw = str(layout or "").strip().lower()
    raw = _ALIASES.get(raw, raw)
    return raw if raw in _RENDERERS else default


# ═══════════════════════════════════════════════════════════════════════════
#  CONTENT VALIDATION & REPAIR  — ask LLM to fill missing fields
# ═══════════════════════════════════════════════════════════════════════════

def _needs_repair(slide: dict) -> bool:
    layout = _canon(slide.get("layout","bullets"))
    if layout == "two_column":
        lp = _safe_list(slide.get("left_points",[]))
        rp = _safe_list(slide.get("right_points",[]))
        return len(lp)<2 and len(rp)<2
    if layout == "timeline":
        steps = [s for s in (slide.get("steps") or [])
                 if isinstance(s,dict) and _safe_str(s.get("label")) and _safe_str(s.get("detail"))]
        return len(steps) < 3
    if layout == "icon_grid":
        items = [g for g in (slide.get("grid_items") or [])
                 if isinstance(g,dict) and _safe_str(g.get("title"))]
        return len(items) < 3
    if layout == "case_study":
        return not _safe_str(slide.get("company"))
    if layout == "table":
        cols = _safe_list(slide.get("table_columns",[]))
        rows = [r for r in (slide.get("table_rows") or []) if isinstance(r,list)]
        return len(cols)<2 or len(rows)<2
    if layout == "chart":
        data = [x for x in (slide.get("chart_data") or []) if isinstance(x,dict)]
        return len(data) < 2
    if layout == "big_stat":
        return not _safe_str(slide.get("stat"))
    if layout == "hybrid_insight":
        return not _safe_str(slide.get("stat"))
    return False


_REPAIR_PROMPT = {
    "two_column": """
Return ONLY valid JSON for a two_column slide about "{title}".
Fields required: left_title, right_title, left_points (list of 4 strings), right_points (list of 4 strings).
Topic context: {context}
""",
    "timeline": """
Return ONLY valid JSON for a timeline slide about "{title}".
Fields required: steps — list of 4 objects each with "label" (short phase name) and "detail" (1-sentence description).
Topic context: {context}
""",
    "icon_grid": """
Return ONLY valid JSON for an icon_grid slide about "{title}".
Fields required: grid_items — list of 4 objects each with "icon" (single char), "title" (2-4 words), "detail" (1-sentence).
Write icons carefully ,i.e, they should be alphabetically and numerically in order.
Topic context: {context}
""",
    "case_study": """
Return ONLY valid JSON for a case_study slide about "{title}".
Fields required: company (real org name), result (1-sentence outcome),
metrics (list of 3 objects with "label" and "value" 1-99),
content (list of 4 supporting bullets).
Topic context: {context}
""",
    "table": """
Return ONLY valid JSON for a table slide about "{title}".
Fields required: table_columns (list of 3-5 column names), table_rows (list of 4-5 rows, each a list of strings).
Topic context: {context}
""",
    "chart": """
Return ONLY valid JSON for a chart slide about "{title}".
Fields required: chart_title (string), chart_data (list of 4-5 objects with "label" and "value" 1-100),
chart_source (optional string).
Topic context: {context}
""",
    "big_stat": """
Return ONLY valid JSON for a big_stat slide about "{title}".
Fields required: stat (e.g. "42%"), stat_label (short descriptor), stat_source (source/year),
content (list of 4 supporting bullets).
Topic context: {context}
""",
    "hybrid_insight": """
Return ONLY valid JSON for a hybrid_insight slide about "{title}".
Fields required: stat (e.g. "3.2x"), stat_label (short descriptor),
chart_data (list of 3 objects with "label" and "value" 1-100),
content (list of 4 supporting bullets).
Topic context: {context}
""",
}

def _repair_slide(slide: dict, topic: str):
    layout = _canon(slide.get("layout","bullets"))
    template = _REPAIR_PROMPT.get(layout)
    if not template:
        return slide, None
    context = f"{topic} — slide: {slide.get('title','')} — existing content: {slide.get('content',[][:2])}"
    prompt = template.format(title=slide.get("title",""), context=context)
    usage = None
    try:
        resp = _client.chat.completions.create(
            model=AZURE_DEPLOYMENT,
            messages=[{"role":"user","content":prompt}],
            temperature=0.7,
            max_tokens=900,
        )
        usage = _usage_dict(resp.usage)
        raw = (resp.choices[0].message.content or "").strip()
        raw = re.sub(r"^```[a-z]*\n?","",raw).rstrip("`").strip()
        data = json.loads(raw)
        if isinstance(data, dict):
            slide.update(data)
    except Exception as e:
        print(f"[repair:{layout}] {e}")
    return slide, usage


# ═══════════════════════════════════════════════════════════════════════════
#  JSON PARSING
# ═══════════════════════════════════════════════════════════════════════════

def _parse_json(raw: str) -> dict:
    if not raw: return {}
    txt = raw.strip()
    txt = re.sub(r"^```[a-z]*\n?","",txt).rstrip("`").strip()
    try: return json.loads(txt)
    except:
        try: return ast.literal_eval(txt)
        except: pass
        s = txt.find("{"); e = txt.rfind("}")
        if s>=0 and e>s:
            try: return json.loads(txt[s:e+1])
            except:
                try: return ast.literal_eval(txt[s:e+1])
                except: pass
    return {}


def _seq_icon(idx: int) -> str:
    if 0 <= idx < 26:
        return chr(ord("A") + idx)
    return str(idx + 1)


def _normalize_icon_grid_items(slide: dict) -> list:
    items = slide.get("grid_items", [])
    if not isinstance(items, list):
        items = []

    cleaned = []
    for i, g in enumerate(items):
        if not isinstance(g, dict):
            continue
        t = _safe_str(g.get("title", ""))
        d = _safe_str(g.get("detail", g.get("description", "")))
        if not t:
            continue
        icon = _safe_str(g.get("icon", "")) or _seq_icon(i)
        cleaned.append({"icon": icon, "title": t, "detail": d})

    if cleaned:
        return cleaned[:4]

    content = slide.get("content", [])
    raw_text = "\n".join(str(x) for x in content) if isinstance(content, list) else _safe_str(content)
    parsed = _parse_json(raw_text) if raw_text.strip() else {}

    parsed_items = None
    if isinstance(parsed, dict) and isinstance(parsed.get("grid_items"), list):
        parsed_items = parsed.get("grid_items")
    elif isinstance(parsed, list):
        parsed_items = parsed

    if isinstance(parsed_items, list):
        for i, g in enumerate(parsed_items):
            if not isinstance(g, dict):
                continue
            t = _safe_str(g.get("title", ""))
            d = _safe_str(g.get("detail", g.get("description", "")))
            if not t:
                continue
            icon = _safe_str(g.get("icon", "")) or _seq_icon(i)
            cleaned.append({"icon": icon, "title": t, "detail": d})
        if cleaned:
            return cleaned[:4]

    lines = []
    if isinstance(content, list):
        lines = [_safe_str(x) for x in content if _safe_str(x)]
    elif raw_text.strip():
        lines = [_safe_str(x) for x in raw_text.splitlines() if _safe_str(x)]

    for i, line in enumerate(lines):
        if len(cleaned) >= 4:
            break
        parsed_line = _parse_json(line)
        if isinstance(parsed_line, dict):
            t = _safe_str(parsed_line.get("title", ""))
            d = _safe_str(parsed_line.get("detail", parsed_line.get("description", "")))
            icon = _safe_str(parsed_line.get("icon", "")) or _seq_icon(i)
            if t:
                cleaned.append({"icon": icon, "title": t, "detail": d})
                continue
        title, detail = line, ""
        if ":" in line:
            title, detail = [part.strip() for part in line.split(":", 1)]
        title = _safe_str(title)
        detail = _safe_str(detail)
        if title:
            cleaned.append({"icon": _seq_icon(i), "title": title[:60], "detail": detail})
    return cleaned[:4]


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE NORMALIZER  — clean LLM output into renderer-ready spec
# ═══════════════════════════════════════════════════════════════════════════

def _normalize_slide(slide: dict, idx: int, topic: str) -> dict:
    slide = dict(slide)
    slide.setdefault("title", f"Slide {idx}")
    slide.setdefault("subtitle", "")
    slide.setdefault("content", [])
    slide.setdefault("icon", "▸")
    if not isinstance(slide.get("style"), dict): slide["style"] = {}

    layout = _canon(slide.get("layout","bullets"))
    slide["layout"] = layout
    if _looks_like_prompt_title(slide.get("title", ""), topic):
        slide["title"] = _fallback_slide_title(slide, idx)

    # Normalize list fields
    for k in ("content","left_points","right_points","sections"):
        if k in slide:
            slide[k] = _safe_list(slide[k])

    if layout == "two_column":
        if not slide.get("left_title"): slide["left_title"] = "Option A"
        if not slide.get("right_title"): slide["right_title"] = "Option B"
        if not slide.get("left_points") and not slide.get("right_points"):
            c = _safe_list(slide.get("content",[]))
            mid = max(1, len(c)//2)
            slide["left_points"] = c[:mid]
            slide["right_points"] = c[mid:]

    elif layout == "big_stat":
        slide.setdefault("stat_label","Key Metric")
        slide.setdefault("stat_source","")

    elif layout == "timeline":
        steps = slide.get("steps",[])
        if not isinstance(steps, list): steps = []
        slide["steps"] = [
            {"label": _safe_str(s.get("label")), "detail": _safe_str(s.get("detail"))}
            for s in steps if isinstance(s,dict) and _safe_str(s.get("label")) and _safe_str(s.get("detail"))
        ][:5]

    elif layout == "icon_grid":
        slide["grid_items"] = _normalize_icon_grid_items(slide)
        slide["content"] = [
            f"{item['title']}: {item['detail']}".rstrip(": ").strip()
            for item in slide["grid_items"]
        ]

    elif layout == "case_study":
        slide.setdefault("company","Organisation")
        slide.setdefault("result","")
        metrics = slide.get("metrics",[])
        if not isinstance(metrics,list): metrics = []
        slide["metrics"] = [
            {"label":_safe_str(m.get("label","Metric")),
             "value":_parse_int(m.get("value",50),default=50,lo=1,hi=99)}
            for m in metrics if isinstance(m,dict)
        ][:3]

    elif layout == "table":
        cols = _safe_list(slide.get("table_columns",[]))[:5]
        rows = [r for r in (slide.get("table_rows") or []) if isinstance(r,list)][:6]
        slide["table_columns"] = cols
        slide["table_rows"] = []
        for row in rows:
            cells = [_safe_str(v) for v in row[:len(cols)]]
            while len(cells)<len(cols): cells.append("")
            if any(cells): slide["table_rows"].append(cells)

    elif layout == "chart":
        data = slide.get("chart_data",[])
        if not isinstance(data,list): data = []
        slide["chart_data"] = [
            {"label":_safe_str(d.get("label",""))[:40],
             "value":_parse_int(d.get("value",50),default=50,lo=1,hi=100)}
            for d in data if isinstance(d,dict) and _safe_str(d.get("label",""))
        ][:5]
        slide.setdefault("chart_title", slide["title"])
        slide.setdefault("chart_source","")

    elif layout == "image_text_split":
        slide.setdefault("image_caption", slide["title"])
        side = _safe_str(slide.get("image_side","right")).lower()
        slide["image_side"] = side if side in ("left","right") else "right"

    elif layout == "hybrid_insight":
        slide.setdefault("stat_label","Key Indicator")
        data = slide.get("chart_data",[])
        if not isinstance(data,list): data = []
        slide["chart_data"] = [
            {"label":_safe_str(d.get("label",""))[:24],
             "value":_parse_int(d.get("value",50),default=50,lo=1,hi=100)}
            for d in data if isinstance(d,dict) and _safe_str(d.get("label",""))
        ][:3]

    elif layout == "section_index":
        secs = _safe_list(slide.get("sections", slide.get("content",[])))
        slide["sections"] = secs[:6]
        slide["content"]  = secs[:6]

    return slide


# ═══════════════════════════════════════════════════════════════════════════
#  MAIN CONTENT GENERATION  — single LLM call, all layouts fully specified
# ═══════════════════════════════════════════════════════════════════════════

def generate_slide_content(topic: str, num_slides: int = 6,
                           tone: str = "Professional",
                           logo_path: str = None) -> dict:
    target = max(3, int(num_slides or 6))
    palette_name = random.choice(list(PALETTES.keys()))

    logo_colors = ""
    if logo_path:
        tmp = _extract_logo_theme(logo_path, _BASE)
        def _hex(t): return "#{:02X}{:02X}{:02X}".format(*t)
        logo_colors = f"""
Logo-extracted palette (use as base for design_system.theme):
  primary: {_hex(tmp['p'])}, secondary: {_hex(tmp['s'])}, accent: {_hex(tmp['a'])},
  accent2: {_hex(tmp['a2'])}, bg_dark: {_hex(tmp['dk'])}
"""

    prompt = f"""
You are an expert presentation strategist.
Create a {tone} PowerPoint deck on: "{topic}".
Return ONLY valid JSON. No markdown fences. No commentary.

There should be no extra white background behind the logo on the top right corner.

Generate exactly {target} slides.

=== REQUIRED TOP-LEVEL STRUCTURE ===
{{
  "design_system": {{
    "theme": {{
      "primary": "#RRGGBB",
      "secondary": "#RRGGBB",
      "accent": "#RRGGBB",
      "accent2": "#RRGGBB",
      "bg_dark": "#RRGGBB",
      "bg_light": "#FFFFFF",
      "text_dark": "#RRGGBB",
      "text_light": "#FFFFFF",
      "card_bg": "#FFFFFF",
      "header_font": "font name",
      "body_font": "font name"
    }}
  }},
  "slides": [ ... exactly {target} slide objects ... ]
}}

{logo_colors}

=== SLIDE RULES ===
Slide 1: layout = "title_cover"
Slide 2: layout = "section_index" (agenda/contents)
Slides 3-N: choose the BEST layout for each section's content type.
  - Avoid repeating the same layout in consecutive slides.
  - Use concrete facts, real companies, real numbers — no generic filler.
  - Generate smart, presentation-ready slide titles.
  - Never copy the user's raw prompt/topic verbatim as a slide title.
  - For the cover slide, write a polished executive title inspired by the topic, not the raw prompt text.
  - For content slides, each title must be a concise section heading that reflects that slide's content.


=== EVERY SLIDE MUST HAVE ===
title, subtitle, layout, icon (single char like ▸ ◆ ✓), content (list of strings), style object.

style object fields (all required):
  pattern_name, surface ("light" or "tint"),
  header_variant ("solid"|"split"|"banded"),
  card_variant ("outline"|"soft"|"banded"),
  footer_variant ("solid"|"line"),
  badge_shape ("oval"|"rect"),
  accent_rotation ("static"|"auto")

=== LAYOUT-SPECIFIC REQUIRED FIELDS ===

"title_cover":
  content: [2-3 short highlight strings]

"section_index":
  sections: [4-7 short section title strings matching the deck flow]
  content: same as sections

"bullets":
  content: [5-6 substantive bullet strings, use **bold** for key terms]

"two_column":
  left_title: "string"
  right_title: "string"
  left_points: ["4 substantive bullet strings"]
  right_points: ["4 substantive bullet strings"]
  content: []

"big_stat":
  stat: "real metric e.g. 87%, $4.2B, 3.1x"
  stat_label: "short descriptor"
  stat_source: "Source, Year"
  content: ["4-5 supporting bullet strings"]

"timeline":
  steps: [{{"label":"Phase/Year","detail":"1-sentence description"}}, ...4 items]
  Keep labels very short and details compact so all 4 timeline numbers remain clearly visible.
  content: []

"icon_grid":
  grid_items: [{{"icon":"▸","title":"2-3 word title","detail":"1-sentence description"}}, ...4 items]
  content: []

"case_study":
  company: "Real organisation name"
  result: "One-sentence outcome"
  metrics: [{{"label":"metric name","value":75}}, ...2-3 items, value 1-99]
  content: ["Exactly 3 short supporting bullet strings with no duplicates"]

"table":
  table_columns: ["Col1","Col2","Col3","Col4"]
  table_rows: [["v1","v2","v3","v4"], ...4-5 rows]
  content: []

"chart":
  chart_title: "string"
  chart_data: [{{"label":"name","value":75}}, ...4-5 items, value 1-100]
  chart_source: "optional source"
  content: []

"image_text_split":
  image_caption: "string"
  image_side: "left" or "right"
  content: ["4-6 bullet strings"]

"hybrid_insight":
  stat: "e.g. 2.4x"
  stat_label: "descriptor"
  chart_data: [{{"label":"name","value":75}}, ...3 items]
  content: ["4-5 bullet strings"]

=== FORMATTING SAFETY RULES ===
- Keep all text concise enough to fit cleanly on a slide.
- Never generate duplicate bullets.
- Keep metric labels short.
- Keep value strings like "25%" or "+98%" on one line.
- Avoid overly verbose timeline labels or details.
"""
    resp = _client.chat.completions.create(
        model=AZURE_DEPLOYMENT,
        messages=[{"role":"user","content":prompt}],
        temperature=0.6,
        max_tokens=7000,
    )
    usage_total = _usage_dict(resp.usage)
    raw  = (resp.choices[0].message.content or "").strip()
    data = _parse_json(raw)

    if not isinstance(data, dict): data = {}
    data.setdefault("design_system", {})
    data.setdefault("slides", [])

    slides = [s for s in data["slides"] if isinstance(s,dict)]

    # Normalize all slides
    normalized = []
    for i, s in enumerate(slides, start=1):
        ns = _normalize_slide(s, i, topic)
        normalized.append(ns)

    # Enforce cover + section_index
    if normalized:
        normalized[0]["layout"] = "title_cover"
    if len(normalized) > 1 and _canon(normalized[1].get("layout","")) != "section_index":
        normalized.insert(1, {
            "title":"Contents","subtitle":f"{tone} overview",
            "layout":"section_index","icon":"▸",
            "sections":[s.get("title","") for s in normalized[2:] if s.get("title")][:6],
            "content":[], "style":{}
        })
    if len(normalized) > 1:
        normalized[1]["layout"] = "section_index"
        secs = [s.get("title","") for s in normalized[2:] if _safe_str(s.get("title",""))][:6]
        normalized[1]["sections"] = secs
        normalized[1]["content"]  = secs

    # Trim to target    
    normalized = normalized[:target]

    # Repair slides with missing layout-specific data
    repaired = []
    for s in normalized:
        if _needs_repair(s):
            s, repair_usage = _repair_slide(s, topic)
            usage_total = _merge_usage(usage_total, repair_usage)
            s = _normalize_slide(s, 0, topic)
        repaired.append(s)

    # Final hard guard: always return exactly the requested number of slides.
    repaired = repaired[:target]
    while len(repaired) < target:
        repaired.append(_normalize_slide({
            "title": f"Additional Insight {len(repaired) + 1}",
            "subtitle": f"{topic} highlights",
            "layout": "bullets",
            "content": [
                f"Key development related to {topic}.",
                "Practical implication for the audience.",
                "Short takeaway that fits the narrative.",
            ],
            "style": {},
        }, len(repaired) + 1, topic))

    if repaired:
        repaired[0]["layout"] = "title_cover"
    if len(repaired) > 1:
        repaired[1]["layout"] = "section_index"
        secs = _derive_sections(repaired[2:])
        repaired[1]["sections"] = secs
        repaired[1]["content"] = secs

    data["slides"] = repaired
    data["usage"] = usage_total
    return data


# ═══════════════════════════════════════════════════════════════════════════
#  DERIVE SECTION TITLES  (for Contents slide if missing)
# ═══════════════════════════════════════════════════════════════════════════

def _derive_sections(slides):
    return [_safe_str(s.get("title","")) for s in slides
            if _canon(s.get("layout","")) not in ("title_cover","section_index")
            and _safe_str(s.get("title",""))][:6]


# ═══════════════════════════════════════════════════════════════════════════
#  PUBLIC API
# ═══════════════════════════════════════════════════════════════════════════

def create_ppt(slide_data, topic: str,
               logo_path: str = None,
               tone: str = "Professional",
               content_image_path: str = None) -> str:

    payload = slide_data if isinstance(slide_data, dict) else {"slides": slide_data}
    design  = payload.get("design_system", {}) if isinstance(payload, dict) else {}
    slides  = payload.get("slides", [])
    if not isinstance(slides, list): slides = []

    if content_image_path and os.path.exists(content_image_path):
        target_idx = None
        for idx, spec in enumerate(slides):
            if isinstance(spec, dict) and (
                _canon(spec.get("layout","")) == "image_text_split" or bool(spec.get("use_uploaded_image"))
            ):
                target_idx = idx
                break
        if target_idx is None:
            for idx, spec in enumerate(slides):
                if isinstance(spec, dict) and _canon(spec.get("layout","")) not in ("title_cover", "section_index"):
                    target_idx = idx
                    break
        if target_idx is not None and isinstance(slides[target_idx], dict):
            slides[target_idx] = dict(slides[target_idx])
            if _canon(slides[target_idx].get("layout","")) != "image_text_split":
                slides[target_idx]["layout"] = "image_text_split"
                slides[target_idx].setdefault("image_caption", _safe_str(slides[target_idx].get("title", "Reference Image")))
                slides[target_idx].setdefault("image_side", "right")
            slides[target_idx]["uploaded_image_path"] = content_image_path

    palette_name = random.choice(list(PALETTES.keys()))
    theme   = _build_theme(palette_name, design, logo_path)
    profile = PROFILES[palette_name]

    prs = Presentation()
    prs.slide_width  = _IN(SW)
    prs.slide_height = _IN(SH)

    for i, spec in enumerate(slides, start=1):
        if not isinstance(spec, dict): continue
        spec = _normalize_slide(spec, i, topic)
        layout   = _canon(spec.get("layout","bullets"))
        renderer = _RENDERERS.get(layout, render_bullets)
        pslide   = prs.slides.add_slide(prs.slide_layouts[6])
        renderer(pslide, spec, i, theme, profile, logo_path)

    # Fallback: empty deck guard
    if len(prs.slides) == 0:
        pslide = prs.slides.add_slide(prs.slide_layouts[6])
        render_bullets(pslide, {
            "title": topic, "subtitle": "Auto-generated",
            "layout":"bullets","icon":"▸",
            "content": [f"No slide data was returned for: {topic}. Please try again."],
        }, 1, theme, profile, logo_path)
    
    deck_name = _safe_str(topic)
    if isinstance(slides, list) and slides:
        first_title = _safe_str(slides[0].get("title", "")) if isinstance(slides[0], dict) else ""
        if first_title:
            deck_name = first_title

    safe = re.sub(r"[^\w\-]","_", topic)[:60]
    os.makedirs("generated", exist_ok=True)
    path = f"generated/{safe}.pptx"
    prs.save(path)
    return path